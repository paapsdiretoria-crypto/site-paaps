#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Converte um deck HTML da PAAPS em PPTX com TEXTO EDITAVEL.

    python3 gerar.py <arquivo.html> <n-slides> <saida.pptx> [--conferir]

A foto, o veu marrom, os paineis e os grafismos viram a imagem de fundo de
cada slide. Todo o texto vira caixa de texto de verdade, na posicao medida no
navegador, com a cor, o corpo, o negrito e o destaque em amarelo preservados.

Quatro armadilhas que este script resolve, todas descobertas na marra:

1. MEDIR E FOTOGRAFAR NA MESMA JANELA. O slide e 16:9 travado. O `--dump-dom`
   do Chrome usa uma janela util mais baixa que o `--screenshot`, e o slide
   muda de largura entre as duas. Por isso a janela e 1920x1167 (que da
   viewport de 1920x1080) e a foto e cortada nos 1080 de cima.

2. CONGELAR A ANIMACAO. Os elementos `.rev` entram com translateY(22px). Em
   duas execucoes separadas do Chrome a animacao pega em estados diferentes.
   `parado.css` forca o estado final.

3. NAO ALTERAR O DOM. Envolver texto solto em <span> mudava o layout flex dos
   cartoes: o fundo saia diferente da peca. `extrair.js` mede por Range e
   marca por atributo, que nao mexe em layout.

4. A LINHA VEM PRONTA. O PowerPoint nao requebra nada: cada linha visual do
   navegador entra como uma linha fixa. E o passo de linha sai da MEDIDA, nao
   da folha de estilo, porque o bloco pode declarar 25px e conter um numero de
   44px dentro de um <span>.

Precisa de: Google Chrome, python-pptx, Pillow. Para `--conferir`, tambem
LibreOffice e pymupdf.
"""
import http.server, json, pathlib, re, shutil, socket, socketserver, subprocess
import sys, threading, html as _html
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

AQUI = pathlib.Path(__file__).resolve().parent
RAIZ = AQUI.parent.parent                      # codigo/site
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"
SOFFICE = "/Applications/LibreOffice.app/Contents/MacOS/soffice"

LARG, ALT = 1920, 1080
JANELA = "1920,1167"        # viewport util de 1920x1080
PX_EMU, PX_PT = 6350.0, 0.5

# Calibragem MEDIDA (ver o cabecalho): o motor do PowerPoint desenha a tinta em
# topo_da_caixa + 0,308 * corpo; o navegador, em topo_do_range + 0,017 * corpo.
# A caixa sobe a diferenca, e a correcao e proporcional ao corpo da letra:
# uma constante acerta o texto corrido e erra o numero grande.
CORRECAO = 0.3079 - 0.0172

ALINHA = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}


# ---------------------------------------------------------------- servidor
class Silencioso(http.server.SimpleHTTPRequestHandler):
    def log_message(self, *a):
        pass


def servir(pasta):
    porta = socket.socket()
    porta.bind(("127.0.0.1", 0))
    n = porta.getsockname()[1]
    porta.close()
    srv = socketserver.TCPServer(
        ("127.0.0.1", n),
        lambda *a, **k: Silencioso(*a, directory=str(pasta), **k))
    threading.Thread(target=srv.serve_forever, daemon=True).start()
    return srv, n


# ------------------------------------------------------------------ copias
def copias(fonte):
    """Gera, ao lado do deck, uma copia congelada com texto e outra sem. Ficam
       no mesmo diretorio porque o CSS e as fotos sao caminhos relativos."""
    h = fonte.read_text()
    h = h.replace("</head>", f'<link rel="stylesheet" href="{AQUI.name}/parado.css">\n</head>', 1)
    fixar = """
<script>
(function(){function f(){document.querySelectorAll('section[class*="slide"]')
.forEach(function(s){s.classList.add('vis')});}f();addEventListener('load',f);})();
</script>
"""
    extrair = f"""
<script src="{AQUI.name}/extrair.js"></script>
<script>
addEventListener('load', function () {{
  setTimeout(function () {{
    var d = window.__pptxEsconder();
    var s = document.createElement('script');
    s.type = 'application/json'; s.id = 'pptxdata';
    s.textContent = JSON.stringify(d);
    document.body.appendChild(s);
  }}, 1200);
}});
</script>
"""
    com = fonte.with_name("_tmp-com-texto.html")
    sem = fonte.with_name("_tmp-sem-texto.html")
    com.write_text(h.replace("</body>", fixar + "</body>", 1))
    sem.write_text(h.replace("</body>", fixar + extrair + "</body>", 1))
    return com, sem


def chrome(args, url):
    subprocess.run([CHROME, "--headless", "--disable-gpu", "--hide-scrollbars",
                    f"--window-size={JANELA}", "--virtual-time-budget=12000",
                    *args, url], capture_output=True)


def capturar(base, sem, com, n, tmp):
    """Para cada slide: a geometria do texto e a foto do fundo sem texto."""
    (tmp / "fundos").mkdir(parents=True, exist_ok=True)
    blocos = {}
    for i in range(1, n + 1):
        dom = tmp / f"dom{i}.html"
        r = subprocess.run([CHROME, "--headless", "--disable-gpu", "--hide-scrollbars",
                            f"--window-size={JANELA}", "--virtual-time-budget=12000",
                            "--dump-dom", f"{base}/{sem}?solo={i}"],
                           capture_output=True, text=True)
        dom.write_text(r.stdout)
        m = re.search(r'id="pptxdata"[^>]*>(.*?)</script>', r.stdout, re.S)
        if not m:
            raise SystemExit(f"slide {i}: o extrator nao devolveu dados")
        o = json.loads(_html.unescape(m.group(1)))
        if tuple(o["janela"][:2]) != (LARG, ALT):
            raise SystemExit(f"janela errada: {o['janela']}. Ajuste JANELA.")
        blocos[i] = o["slides"][i - 1]["itens"]

        png = tmp / f"f{i}.png"
        chrome(["--screenshot=" + str(png)], f"{base}/{sem}?solo={i}")
        im = Image.open(png).convert("RGB").crop((0, 0, LARG, ALT))
        im.save(tmp / "fundos" / f"s{i}.jpg", "JPEG", quality=92,
                subsampling=0, optimize=True)
        if com:
            ref = tmp / f"ref{i}.png"
            chrome(["--screenshot=" + str(ref)], f"{base}/{com}?solo={i}")
            Image.open(ref).convert("RGB").crop((0, 0, LARG, ALT)).save(
                tmp / f"ref{i}.png")
        print(f"  slide {i}/{n}")
    return blocos


# -------------------------------------------------------------------- pptx
def cor(t):
    v = [float(x) for x in re.match(r"rgba?\(([^)]+)\)", t).group(1).split(",")]
    return (v[0], v[1], v[2], v[3] if len(v) > 3 else 1.0)


def sobre(rgba, chao):
    """O PowerPoint nao tem alfa de texto pela via simples: a cor e composta
       sobre o pixel medio do fundo, que e o que o olho ve de qualquer jeito."""
    r, g, b, a = rgba
    return RGBColor(*(int(round(v * a + c * (1 - a))) for v, c in zip((r, g, b), chao)))


def media(im, x, y, w, h):
    x0, y0 = max(0, int(x)), max(0, int(y))
    x1, y1 = min(im.width, int(x + w) + 1), min(im.height, int(y + h) + 1)
    if x1 <= x0 or y1 <= y0:
        return (30, 14, 3)
    return im.crop((x0, y0, x1, y1)).resize((1, 1), Image.LANCZOS).getpixel((0, 0))[:3]


def corpo(l):
    return max(r["s"] for r in l["runs"])


def agrupar(lns, lh):
    """Quebra o bloco em grupos de linhas REGULARES, com o passo MEDIDO.
       Um rotulo pequeno em cima de um numero grande sao corpos diferentes:
       reespacar os dois pela mesma entrelinha faz um cair sobre o outro."""
    grupos, atual = [], []
    for l in lns:
        if not atual:
            atual = [l]
            continue
        ant = atual[-1]
        if abs(corpo(l) - corpo(ant)) < 0.6 and abs((l["y"] - ant["y"]) - lh) <= 1.6:
            atual.append(l)
        else:
            grupos.append(atual)
            atual = [l]
    if atual:
        grupos.append(atual)

    saida = []
    for g in grupos:
        if len(g) > 1:
            passo = (g[-1]["y"] - g[0]["y"]) / (len(g) - 1)
        else:
            passo = max(g[0]["h"], corpo(g[0]) * 1.25)
        saida.append((g, passo))
    return saida


def caixa(slide, im, b, lns, passo):
    al = b["align"]
    esq = min(l["x"] for l in lns)
    dirr = max(l["x"] + l["w"] for l in lns)

    # A caixa precisa ser mais larga que a linha, senao a menor diferenca de
    # metrica requebra o texto. A folga vai toda no lado que NAO ancora.
    folga = max(60.0, (dirr - esq) * 0.18)
    if al == "right":
        x, larg = esq - folga, (dirr - esq) + folga
    elif al == "center":
        x, larg = esq - folga / 2, (dirr - esq) + folga
    else:
        x, larg = esq, (dirr - esq) + folga
    if x < 0:
        larg += x
        x = 0.0
    larg = min(larg, LARG - x)

    y = min(l["y"] for l in lns) - CORRECAO * corpo(lns[0])
    cx = slide.shapes.add_textbox(Emu(int(x * PX_EMU)), Emu(int(y * PX_EMU)),
                                  Emu(int(larg * PX_EMU)),
                                  Emu(int(passo * len(lns) * PX_EMU)))
    tf = cx.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for k, l in enumerate(lns):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = ALINHA.get(al, PP_ALIGN.LEFT)
        p.line_spacing = Pt(round(passo * PX_PT, 2))
        p.space_before = p.space_after = Pt(0)
        chao = media(im, l["x"], l["y"], l["w"], l["h"])
        for t in l["runs"]:
            r = p.add_run()
            r.text = t["t"]
            f = r.font
            f.name = t["f"]
            f.size = Pt(round(t["s"] * PX_PT, 1))
            f.bold = t["b"]
            f.italic = t["i"]
            f.color.rgb = sobre(cor(t["c"]), chao)
            if abs(b["ls"]) >= 0.1:
                f._rPr.set("spc", str(int(round(b["ls"] * PX_PT * 100))))


def montar(blocos, tmp, saida, n):
    prs = Presentation()
    prs.slide_width = Emu(int(LARG * PX_EMU))
    prs.slide_height = Emu(int(ALT * PX_EMU))
    vazio = prs.slide_layouts[6]
    for i in range(1, n + 1):
        s = prs.slides.add_slide(vazio)
        fundo = tmp / "fundos" / f"s{i}.jpg"
        s.shapes.add_picture(str(fundo), 0, 0, prs.slide_width, prs.slide_height)
        im = Image.open(fundo).convert("RGB")
        for b in blocos[i]:
            for g, passo in agrupar(b["linhas"], b["lh"]):
                caixa(s, im, b, g, passo)
    prs.save(saida)


def conferir(saida, tmp, n):
    """Renderiza o PPTX e compara com o navegador, slide a slide."""
    import numpy as np
    import pymupdf
    out = tmp / "render"
    out.mkdir(exist_ok=True)
    subprocess.run([SOFFICE, "--headless", "--convert-to", "pdf",
                    "--outdir", str(out), str(saida)], capture_output=True)
    d = pymupdf.open(out / (saida.stem + ".pdf"))
    print("\n  slide | desvio (dx,dy) | tinta em comum")
    for i in range(1, n + 1):
        pg = d[i - 1]
        z = LARG / pg.rect.width
        pg.get_pixmap(matrix=pymupdf.Matrix(z, z)).save(out / f"r{i}.png")
        ler = lambda p: np.asarray(Image.open(p).convert("L").resize((LARG, ALT)),
                                   dtype=np.float32)
        f = ler(tmp / "fundos" / f"s{i}.jpg")
        mo = np.abs(ler(tmp / f"ref{i}.png") - f) > 30
        mr = np.abs(ler(out / f"r{i}.png") - f) > 30
        melhor = (0, 0, -1.0)
        for dy in range(-8, 9, 2):
            for dx in range(-8, 9, 2):
                v = float((mo & np.roll(np.roll(mr, -dy, 0), -dx, 1)).sum())
                if v > melhor[2]:
                    melhor = (dx, dy, v)
        dx, dy, inter = melhor
        uni = float((mo | np.roll(np.roll(mr, -dy, 0), -dx, 1)).sum())
        print(f"   {i:>3}  |   {dx:+3d},{dy:+3d}      |  {100*inter/max(uni,1):5.1f}%")


def main():
    if len(sys.argv) < 4:
        raise SystemExit(__doc__)
    fonte = pathlib.Path(sys.argv[1]).resolve()
    n = int(sys.argv[2])
    saida = pathlib.Path(sys.argv[3]).resolve()
    checar = "--conferir" in sys.argv

    if not pathlib.Path(CHROME).exists():
        raise SystemExit(f"Chrome nao encontrado em {CHROME}")

    tmp = saida.parent / f".tmp-{saida.stem}"
    if tmp.exists():
        shutil.rmtree(tmp)
    tmp.mkdir(parents=True)

    com, sem = copias(fonte)
    srv, porta = servir(RAIZ)
    try:
        rel = fonte.parent.relative_to(RAIZ)
        base = f"http://127.0.0.1:{porta}/{rel}"
        print(f"Lendo {fonte.name} em {n} slides...")
        blocos = capturar(base, sem.name, com.name if checar else None, n, tmp)
        print("Montando o PPTX...")
        montar(blocos, tmp, saida, n)
        print(f"Pronto: {saida}  ({saida.stat().st_size/1e6:.1f} MB, "
              f"{sum(len(v) for v in blocos.values())} blocos de texto)")
        if checar:
            conferir(saida, tmp, n)
    finally:
        srv.shutdown()
        for f in (com, sem):
            f.unlink(missing_ok=True)
        if not checar:
            shutil.rmtree(tmp, ignore_errors=True)


if __name__ == "__main__":
    main()
