# -*- coding: utf-8 -*-
"""Deck-sonda: caixas de texto de corpos conhecidos em posicoes conhecidas.
   Renderiza, mede onde a tinta caiu de verdade e calcula a correcao. E o
   unico jeito honesto de saber como o motor posiciona a primeira linha."""
import subprocess, pathlib, sys
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

PX_EMU, PX_PT = 6350.0, 0.5
CORPOS = [9, 12, 16, 20, 26, 34, 44, 56, 70]
TOPO, PASSO_LINHA = 60.0, 110.0

prs = Presentation()
prs.slide_width = Emu(int(1920 * PX_EMU)); prs.slide_height = Emu(int(1080 * PX_EMU))
s = prs.slides.add_slide(prs.slide_layouts[6])
fundo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
fundo.fill.solid(); fundo.fill.fore_color.rgb = RGBColor(0, 0, 0); fundo.line.fill.background()

for i, c in enumerate(CORPOS):
    passo = c * 1.25
    topo = TOPO + i * PASSO_LINHA
    cx = s.shapes.add_textbox(Emu(int(120 * PX_EMU)), Emu(int(topo * PX_EMU)),
                              Emu(int(900 * PX_EMU)), Emu(int(passo * PX_EMU)))
    tf = cx.text_frame
    tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    p.line_spacing = Pt(round(passo * PX_PT, 2))
    r = p.add_run(); r.text = "HXhxg 0123"
    r.font.name = "Helvetica"; r.font.size = Pt(round(c * PX_PT, 1))
    r.font.bold = True; r.font.color.rgb = RGBColor(255, 255, 255)

prs.save("sonda.pptx")
subprocess.run(["/Applications/LibreOffice.app/Contents/MacOS/soffice", "--headless",
                "--convert-to", "pdf", "--outdir", "sonda-out", "sonda.pptx"],
               capture_output=True)
import pymupdf
d = pymupdf.open("sonda-out/sonda.pdf"); pg = d[0]
z = 1920 / pg.rect.width
pg.get_pixmap(matrix=pymupdf.Matrix(z, z)).save("sonda.png")

im = np.asarray(Image.open("sonda.png").convert("L").resize((1920, 1080)), dtype=np.float32)
print(" corpo | topo da caixa | topo da tinta | delta | delta/corpo")
ks = []
for i, c in enumerate(CORPOS):
    topo = TOPO + i * PASSO_LINHA
    faixa = im[int(topo) - 30: int(topo) + int(PASSO_LINHA) - 10, :]
    ys, xs = np.nonzero(faixa > 90)
    if not len(ys):
        print(f" {c:>5} | sem tinta"); continue
    tinta = int(topo) - 30 + int(ys.min())
    d_ = tinta - topo
    ks.append(d_ / c)
    print(f" {c:>5} | {topo:13.1f} | {tinta:13d} | {d_:+5.1f} | {d_/c:+.3f}")
print("\nk medio (delta = k * corpo):", round(float(np.mean(ks)), 4))
